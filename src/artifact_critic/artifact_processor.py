"""
Artifact preprocessing utilities for the Artifact Critic.
Handles PDF, PPTX, and image processing.
"""
import io
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from PIL import Image
import base64

# Try to import PDF/PPTX libraries
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("Warning: PyPDF2 not available. PDF text extraction disabled.")

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    print("Warning: pdf2image not available. PDF to image conversion disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. PPTX processing disabled.")


@dataclass
class ProcessedArtifact:
    """Processed artifact ready for review."""
    artifact_type: str
    artifact_id: str
    total_pages: int
    images: List[bytes]  # Image data as bytes
    page_metadata: List[Dict[str, Any]]  # Metadata per page/slide
    extracted_text: Optional[str] = None
    file_metadata: Optional[Dict[str, Any]] = None


class ArtifactProcessor:
    """Processes various artifact types into standardized format."""
    
    def __init__(self, temp_dir: str = "./temp_artifacts", max_pages: int = 10):
        self.temp_dir = Path(temp_dir)
        self.temp_dir.mkdir(parents=True, exist_ok=True)
        self.max_pages = max_pages
    
    def process_artifact(
        self,
        file_path: str,
        artifact_type: Optional[str] = None,
        high_detail: bool = False
    ) -> ProcessedArtifact:
        """
        Process an artifact file into standardized format.
        
        Args:
            file_path: Path to artifact file
            artifact_type: Optional type hint (pdf, pptx, image)
            high_detail: Use higher resolution for text-heavy artifacts
            
        Returns:
            ProcessedArtifact object
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"Artifact not found: {file_path}")
        
        # Detect type if not provided
        if artifact_type is None:
            artifact_type = self._detect_type(file_path)
        
        # Route to appropriate processor
        artifact_type = artifact_type.lower()
        
        if artifact_type == "pdf":
            return self._process_pdf(file_path, high_detail)
        elif artifact_type in ["pptx", "ppt"]:
            return self._process_pptx(file_path, high_detail)
        elif artifact_type in ["png", "jpg", "jpeg", "image"]:
            return self._process_image(file_path, high_detail)
        else:
            raise ValueError(f"Unsupported artifact type: {artifact_type}")
    
    def _detect_type(self, file_path: Path) -> str:
        """Detect artifact type from file extension."""
        ext = file_path.suffix.lower()
        
        if ext == ".pdf":
            return "pdf"
        elif ext in [".pptx", ".ppt"]:
            return "pptx"
        elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp"]:
            return "image"
        else:
            raise ValueError(f"Cannot detect type for: {ext}")
    
    def _process_pdf(self, file_path: Path, high_detail: bool) -> ProcessedArtifact:
        """Process PDF file."""
        if not PDF2IMAGE_AVAILABLE:
            raise RuntimeError("pdf2image not available. Install with: pip install pdf2image")
        
        # Convert PDF to images
        dpi = 300 if high_detail else 150
        
        try:
            images = convert_from_path(
                str(file_path),
                dpi=dpi,
                fmt='png'
            )
        except Exception as e:
            raise RuntimeError(f"Failed to convert PDF to images: {e}")
        
        # Limit pages
        if len(images) > self.max_pages:
            print(f"Warning: PDF has {len(images)} pages, limiting to {self.max_pages}")
            images = images[:self.max_pages]
        
        # Convert to bytes
        image_bytes = []
        page_metadata = []
        
        for i, img in enumerate(images):
            # Convert PIL Image to bytes
            buf = io.BytesIO()
            img.save(buf, format='PNG')
            image_bytes.append(buf.getvalue())
            
            page_metadata.append({
                "page": i + 1,
                "width": img.width,
                "height": img.height,
                "format": "png"
            })
        
        # Extract text if available
        extracted_text = None
        if PDF_AVAILABLE:
            try:
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    text_parts = []
                    for i, page in enumerate(pdf_reader.pages[:self.max_pages]):
                        text_parts.append(f"--- Page {i+1} ---")
                        text_parts.append(page.extract_text())
                    extracted_text = "\n".join(text_parts)
            except Exception:
                pass
        
        return ProcessedArtifact(
            artifact_type="pdf",
            artifact_id=file_path.stem,
            total_pages=len(image_bytes),
            images=image_bytes,
            page_metadata=page_metadata,
            extracted_text=extracted_text,
            file_metadata={
                "filename": file_path.name,
                "size_bytes": file_path.stat().st_size,
                "dpi": dpi
            }
        )
    
    def _process_pptx(self, file_path: Path, high_detail: bool) -> ProcessedArtifact:
        """Process PowerPoint file."""
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not available. Install with: pip install python-pptx")
        
        # Note: Converting PPTX to images requires additional tools
        # For now, we'll extract text and provide a placeholder
        
        try:
            prs = Presentation(str(file_path))
            
            text_parts = []
            page_metadata = []
            
            for i, slide in enumerate(prs.slides[:self.max_pages]):
                text_parts.append(f"--- Slide {i+1} ---")
                
                # Extract text from shapes
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                text_parts.append("\n".join(slide_text))
                
                page_metadata.append({
                    "slide": i + 1,
                    "shapes": len(slide.shapes)
                })
            
            extracted_text = "\n".join(text_parts)
            
            # For now, create placeholder images
            # In production, you'd use a tool like LibreOffice/unoconv to convert
            image_bytes = []
            print("Warning: PPTX to image conversion not fully implemented. Using text extraction only.")
            
            return ProcessedArtifact(
                artifact_type="pptx",
                artifact_id=file_path.stem,
                total_pages=len(page_metadata),
                images=image_bytes,
                page_metadata=page_metadata,
                extracted_text=extracted_text,
                file_metadata={
                    "filename": file_path.name,
                    "size_bytes": file_path.stat().st_size,
                    "total_slides": len(prs.slides)
                }
            )
            
        except Exception as e:
            raise RuntimeError(f"Failed to process PPTX: {e}")
    
    def _process_image(self, file_path: Path, high_detail: bool) -> ProcessedArtifact:
        """Process image file."""
        try:
            img = Image.open(file_path)
            
            # Resize if needed for high detail
            if high_detail and max(img.size) < 1920:
                # Upscale smaller images
                scale = 1920 / max(img.size)
                new_size = (int(img.width * scale), int(img.height * scale))
                img = img.resize(new_size, Image.LANCZOS)
            elif not high_detail and max(img.size) > 1920:
                # Downscale large images
                scale = 1920 / max(img.size)
                new_size = (int(img.width * scale), int(img.height * scale))
                img = img.resize(new_size, Image.LANCZOS)
            
            # Convert to bytes
            buf = io.BytesIO()
            img.save(buf, format='PNG')
            image_bytes = [buf.getvalue()]
            
            return ProcessedArtifact(
                artifact_type="image",
                artifact_id=file_path.stem,
                total_pages=1,
                images=image_bytes,
                page_metadata=[{
                    "page": 1,
                    "width": img.width,
                    "height": img.height,
                    "format": img.format
                }],
                extracted_text=None,
                file_metadata={
                    "filename": file_path.name,
                    "size_bytes": file_path.stat().st_size
                }
            )
            
        except Exception as e:
            raise RuntimeError(f"Failed to process image: {e}")
    
    def encode_image_base64(self, image_bytes: bytes) -> str:
        """Encode image bytes to base64 string for API."""
        return base64.b64encode(image_bytes).decode('utf-8')
